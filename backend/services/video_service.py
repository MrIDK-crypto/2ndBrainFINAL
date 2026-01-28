"""
Video Generation Service
AI-powered video generation from documents and knowledge gaps.
Uses Azure TTS for professional voice quality.
Supports Gamma API for AI-powered presentation generation.
"""

import os
import io
import json
import tempfile
import threading
import queue
from datetime import datetime, timezone
from typing import List, Dict, Optional, Tuple, Any, Callable
from dataclasses import dataclass
from pathlib import Path
from enum import Enum
import time

from sqlalchemy.orm import Session

from database.models import (
    Video, Document, KnowledgeGap, GapAnswer, Project, Tenant,
    VideoStatus, DocumentClassification,
    generate_uuid, utc_now
)

# Import Gamma service
try:
    from services.gamma_service import get_gamma_service
    GAMMA_AVAILABLE = True
except ImportError:
    GAMMA_AVAILABLE = False

# Optional imports for video generation
try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    from moviepy.editor import (
        ImageClip, AudioFileClip, concatenate_videoclips,
        CompositeVideoClip, TextClip
    )
    MOVIEPY_AVAILABLE = True
except ImportError:
    MOVIEPY_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# Azure TTS Configuration
AZURE_TTS_KEY = os.getenv("AZURE_TTS_KEY", os.getenv("AZURE_OPENAI_API_KEY"))
AZURE_TTS_REGION = os.getenv("AZURE_TTS_REGION", "eastus2")
AZURE_TTS_VOICE = os.getenv("AZURE_TTS_VOICE", "en-US-JennyNeural")


@dataclass
class SlideContent:
    """Content for a single slide"""
    title: str
    content: str  # Bullet points or text
    notes: str  # Speaker notes / narration text
    image_path: Optional[str] = None
    duration_hint: float = 10.0  # Suggested duration in seconds


@dataclass
class VideoProgress:
    """Video generation progress"""
    status: str
    progress_percent: int
    current_step: str
    error: Optional[str] = None


class VideoService:
    """
    Video generation service.

    Features:
    - Generate videos from documents or knowledge base
    - AI-generated scripts from content
    - Azure TTS for professional voice
    - Progress tracking
    - Background processing
    """

    # Video settings
    VIDEO_WIDTH = 1920
    VIDEO_HEIGHT = 1080
    FPS = 24
    FONT_TITLE_SIZE = 72
    FONT_CONTENT_SIZE = 48
    FONT_PATH = None  # Will be set based on OS

    # Colors
    BG_COLOR = (15, 23, 42)  # Dark blue (Tailwind slate-900)
    TITLE_COLOR = (255, 255, 255)  # White
    CONTENT_COLOR = (226, 232, 240)  # Light gray (Tailwind slate-200)
    ACCENT_COLOR = (59, 130, 246)  # Blue (Tailwind blue-500)

    def __init__(self, db: Session):
        self.db = db
        self._progress_callbacks: Dict[str, Callable] = {}
        self._set_font_path()

    def _set_font_path(self):
        """Set font path based on OS"""
        import platform
        system = platform.system()

        if system == "Darwin":  # macOS
            font_paths = [
                "/System/Library/Fonts/SFNSDisplay.ttf",
                "/System/Library/Fonts/Helvetica.ttc",
                "/Library/Fonts/Arial.ttf"
            ]
        elif system == "Windows":
            font_paths = [
                "C:/Windows/Fonts/segoeui.ttf",
                "C:/Windows/Fonts/arial.ttf"
            ]
        else:  # Linux
            font_paths = [
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/usr/share/fonts/TTF/DejaVuSans.ttf"
            ]

        for path in font_paths:
            if Path(path).exists():
                self.FONT_PATH = path
                break

    # ========================================================================
    # VIDEO CREATION
    # ========================================================================

    def create_video_from_documents(
        self,
        tenant_id: str,
        title: str,
        document_ids: List[str],
        description: Optional[str] = None,
        project_id: Optional[str] = None,
        use_gamma: bool = True  # NEW: Use Gamma API by default
    ) -> Tuple[Optional[Video], Optional[str]]:
        """
        Create a training video from documents.

        Args:
            tenant_id: Tenant ID
            title: Video title
            document_ids: Documents to include
            description: Video description
            project_id: Optional project association
            use_gamma: Use Gamma API for presentation generation (default True)

        Returns:
            (Video, error)
        """
        try:
            # Create video record
            video = Video(
                tenant_id=tenant_id,
                project_id=project_id,
                title=title,
                description=description,
                status=VideoStatus.QUEUED,
                source_type="documents",
                source_document_ids=document_ids,
                source_config={"use_gamma": use_gamma and GAMMA_AVAILABLE}
            )
            self.db.add(video)
            self.db.commit()

            # Start background processing
            thread = threading.Thread(
                target=self._process_video,
                args=(video.id, tenant_id)
            )
            thread.start()

            return video, None

        except Exception as e:
            self.db.rollback()
            return None, str(e)

    def create_video_from_gaps(
        self,
        tenant_id: str,
        title: str,
        gap_ids: List[str],
        include_answers: bool = True,
        description: Optional[str] = None,
        use_gamma: bool = True  # NEW: Use Gamma API by default
    ) -> Tuple[Optional[Video], Optional[str]]:
        """
        Create a training video from knowledge gaps and answers.

        Args:
            tenant_id: Tenant ID
            title: Video title
            gap_ids: Knowledge gap IDs
            include_answers: Include answers in video
            description: Video description
            use_gamma: Use Gamma API for presentation generation (default True)

        Returns:
            (Video, error)
        """
        try:
            video = Video(
                tenant_id=tenant_id,
                title=title,
                description=description,
                status=VideoStatus.QUEUED,
                source_type="knowledge_gaps",
                source_document_ids=gap_ids,
                source_config={
                    "include_answers": include_answers,
                    "use_gamma": use_gamma and GAMMA_AVAILABLE
                }
            )
            self.db.add(video)
            self.db.commit()

            thread = threading.Thread(
                target=self._process_video,
                args=(video.id, tenant_id)
            )
            thread.start()

            return video, None

        except Exception as e:
            self.db.rollback()
            return None, str(e)

    def _process_video(self, video_id: str, tenant_id: str):
        """Background video processing"""
        # Get new db session for thread
        from database.models import SessionLocal
        db = SessionLocal()

        try:
            video = db.query(Video).filter(Video.id == video_id).first()
            if not video:
                return

            video.status = VideoStatus.PROCESSING
            video.started_at = utc_now()
            video.progress_percent = 0
            db.commit()

            # Get tenant for output path
            tenant = db.query(Tenant).filter(Tenant.id == tenant_id).first()
            if not tenant or not tenant.data_directory:
                raise Exception("Tenant data directory not configured")

            output_dir = Path(tenant.data_directory) / "videos"
            output_dir.mkdir(parents=True, exist_ok=True)

            # Check if we should use Gamma API
            use_gamma = video.source_config.get("use_gamma", False) and GAMMA_AVAILABLE

            if use_gamma:
                # Use Gamma API for presentation generation
                self._update_progress(db, video, 10, "Generating presentation with Gamma AI...")
                slides = self._generate_slides_with_gamma(
                    db, video, tenant_id, output_dir
                )
            else:
                # Generate slides locally based on source type
                self._update_progress(db, video, 10, "Generating content...")

                if video.source_type == "documents":
                    slides = self._generate_slides_from_documents(
                        db, video.source_document_ids, tenant_id
                    )
                elif video.source_type == "knowledge_gaps":
                    slides = self._generate_slides_from_gaps(
                        db,
                        video.source_document_ids,
                        tenant_id,
                        video.source_config.get("include_answers", True)
                    )
                else:
                    slides = []

            if not slides:
                raise Exception("No content to generate video from")

            video.slides_count = len(slides)
            db.commit()

            # Generate audio for each slide
            self._update_progress(db, video, 30, "Generating narration...")
            audio_files = self._generate_audio(slides, output_dir)

            # Render slides to images
            self._update_progress(db, video, 50, "Rendering slides...")
            slide_images = self._render_slides(slides, output_dir)

            # Combine into video
            self._update_progress(db, video, 70, "Creating video...")
            output_path = output_dir / f"{video.id}.mp4"

            duration = self._create_video(
                slide_images,
                audio_files,
                str(output_path)
            )

            # Generate thumbnail
            self._update_progress(db, video, 90, "Generating thumbnail...")
            thumbnail_path = output_dir / f"{video.id}_thumb.jpg"
            self._generate_thumbnail(slide_images[0], str(thumbnail_path))

            # Update video record
            video.status = VideoStatus.COMPLETED
            video.completed_at = utc_now()
            video.progress_percent = 100
            video.file_path = str(output_path)
            video.thumbnail_path = str(thumbnail_path)
            video.duration_seconds = duration
            video.file_size_bytes = output_path.stat().st_size

            db.commit()

            # Cleanup temp files
            for f in audio_files + slide_images:
                try:
                    if f and Path(f).exists():
                        Path(f).unlink()
                except Exception:
                    pass

        except Exception as e:
            video = db.query(Video).filter(Video.id == video_id).first()
            if video:
                video.status = VideoStatus.FAILED
                video.error_message = str(e)
                video.progress_percent = 0
                db.commit()

        finally:
            db.close()

    def _update_progress(
        self,
        db: Session,
        video: Video,
        percent: int,
        step: str
    ):
        """Update video progress"""
        video.progress_percent = percent
        db.commit()

        # Call progress callback if registered
        if video.id in self._progress_callbacks:
            self._progress_callbacks[video.id](VideoProgress(
                status=video.status.value,
                progress_percent=percent,
                current_step=step
            ))

    # ========================================================================
    # CONTENT GENERATION
    # ========================================================================

    def _generate_slides_from_documents(
        self,
        db: Session,
        document_ids: List[str],
        tenant_id: str
    ) -> List[SlideContent]:
        """Generate slide content from documents"""
        documents = db.query(Document).filter(
            Document.id.in_(document_ids),
            Document.tenant_id == tenant_id
        ).all()

        if not documents:
            return []

        # Use GPT to generate slide content
        from services.openai_client import get_openai_client

        client = get_openai_client()

        # Build document summary
        doc_texts = []
        for doc in documents:
            doc_texts.append(f"""
Title: {doc.title or 'Untitled'}
Content: {(doc.content or '')[:2000]}
""")

        combined = "\n---\n".join(doc_texts)

        prompt = f"""Create a training presentation from the following documents.
Generate 5-10 slides that cover the key points.

DOCUMENTS:
{combined}

For each slide provide:
1. Title (short, max 8 words)
2. Content (3-5 bullet points)
3. Narration script (2-3 sentences that a narrator would speak)

Respond in JSON format:
{{
    "slides": [
        {{
            "title": "...",
            "content": "• Point 1\\n• Point 2\\n• Point 3",
            "narration": "..."
        }}
    ]
}}
"""

        try:
            response = client.chat_completion(
                messages=[
                    {"role": "system", "content": "You are a training content creator. Create clear, concise slides."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,
                max_tokens=2000,
                response_format={"type": "json_object"}
            )

            result = json.loads(response.choices[0].message.content)

            slides = []
            for slide_data in result.get("slides", []):
                slides.append(SlideContent(
                    title=slide_data.get("title", ""),
                    content=slide_data.get("content", ""),
                    notes=slide_data.get("narration", "")
                ))

            return slides

        except Exception as e:
            print(f"Error generating slides: {e}")
            return []

    def _generate_slides_from_gaps(
        self,
        db: Session,
        gap_ids: List[str],
        tenant_id: str,
        include_answers: bool
    ) -> List[SlideContent]:
        """Generate slides from knowledge gaps"""
        gaps = db.query(KnowledgeGap).filter(
            KnowledgeGap.id.in_(gap_ids),
            KnowledgeGap.tenant_id == tenant_id
        ).all()

        slides = []

        # Title slide
        slides.append(SlideContent(
            title="Knowledge Transfer",
            content="Key information and answers\nfrom your organization",
            notes="Welcome to this knowledge transfer session. We'll cover important questions and answers from your team."
        ))

        for gap in gaps:
            # Gap intro slide
            slides.append(SlideContent(
                title=gap.title,
                content=gap.description or "",
                notes=f"Let's discuss {gap.title}. {gap.description or ''}"
            ))

            if include_answers:
                # Get answers for this gap
                answers = db.query(GapAnswer).filter(
                    GapAnswer.knowledge_gap_id == gap.id
                ).all()

                for answer in answers:
                    slides.append(SlideContent(
                        title=f"Q: {answer.question_text[:50]}...",
                        content=answer.answer_text[:500],
                        notes=f"Question: {answer.question_text}. Answer: {answer.answer_text}"
                    ))

        # Closing slide
        slides.append(SlideContent(
            title="Summary",
            content="Thank you for watching!\n\nFor more information, consult the knowledge base.",
            notes="That concludes this knowledge transfer session. Thank you for watching."
        ))

        return slides

    def _generate_slides_with_gamma(
        self,
        db: Session,
        video: Video,
        tenant_id: str,
        output_dir: Path
    ) -> List[SlideContent]:
        """
        Generate slides using Gamma API.

        Workflow:
        1. Get source documents/gaps
        2. Generate Gamma presentation
        3. Export as PPTX
        4. Download PPTX
        5. Parse PPTX to extract slide content

        Args:
            db: Database session
            video: Video object
            tenant_id: Tenant ID
            output_dir: Output directory for temp files

        Returns:
            List of SlideContent objects
        """
        try:
            gamma_service = get_gamma_service()

            # Prepare content based on source type
            if video.source_type == "documents":
                documents = db.query(Document).filter(
                    Document.id.in_(video.source_document_ids),
                    Document.tenant_id == tenant_id
                ).all()

                content = gamma_service.generate_from_documents(
                    documents=documents,
                    title=video.title
                )

            elif video.source_type == "knowledge_gaps":
                gaps = db.query(KnowledgeGap).filter(
                    KnowledgeGap.id.in_(video.source_document_ids),
                    KnowledgeGap.tenant_id == tenant_id
                ).all()

                # Get answers if included
                include_answers = video.source_config.get("include_answers", True)
                answers = []
                if include_answers:
                    gap_ids = [g.id for g in gaps]
                    answers = db.query(GapAnswer).filter(
                        GapAnswer.knowledge_gap_id.in_(gap_ids)
                    ).all()

                content = gamma_service.generate_from_knowledge_gaps(
                    gaps=gaps,
                    answers=answers,
                    title=video.title
                )
            else:
                raise Exception(f"Unsupported source type: {video.source_type}")

            print(f"[VideoService] Gamma content prepared ({len(content)} chars)")

            # Generate presentation with PPTX export
            result, error = gamma_service.generate_presentation(
                content=content,
                title=video.title,
                export_format='pptx'
            )

            if error:
                raise Exception(f"Gamma generation failed: {error}")

            if not result or 'exportUrl' not in result:
                raise Exception("Gamma did not return export URL")

            export_url = result['exportUrl']
            print(f"[VideoService] Gamma presentation created, downloading PPTX...")

            # Download PPTX
            pptx_path = output_dir / f"{video.id}_gamma.pptx"
            success, download_error = gamma_service.download_export(
                export_url=export_url,
                output_path=str(pptx_path)
            )

            if not success:
                raise Exception(f"PPTX download failed: {download_error}")

            print(f"[VideoService] PPTX downloaded, parsing slides...")

            # Parse PPTX to extract slides
            slides = self._parse_pptx_to_slides(pptx_path)

            print(f"[VideoService] Extracted {len(slides)} slides from Gamma PPTX")

            # Store Gamma URL in video metadata
            if 'url' in result:
                video.source_config['gamma_url'] = result['url']
                db.commit()

            return slides

        except Exception as e:
            print(f"[VideoService] Gamma generation error: {e}")
            # Fallback to local generation
            print(f"[VideoService] Falling back to local slide generation...")

            if video.source_type == "documents":
                return self._generate_slides_from_documents(
                    db, video.source_document_ids, tenant_id
                )
            else:
                return self._generate_slides_from_gaps(
                    db,
                    video.source_document_ids,
                    tenant_id,
                    video.source_config.get("include_answers", True)
                )

    def _parse_pptx_to_slides(self, pptx_path: Path) -> List[SlideContent]:
        """
        Parse PPTX file to extract slide content.

        Args:
            pptx_path: Path to PPTX file

        Returns:
            List of SlideContent objects
        """
        if not PPTX_AVAILABLE:
            raise Exception("python-pptx library not available")

        slides = []
        prs = Presentation(str(pptx_path))

        for slide_num, slide in enumerate(prs.slides):
            # Extract title
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()

            # Extract content (bullet points and text)
            content_parts = []
            notes_parts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()

                    # Skip title (already extracted)
                    if text == title:
                        continue

                    # Add to content
                    content_parts.append(text)

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        notes_parts.append(notes_text)

            # Build final content and notes
            content = '\n\n'.join(content_parts)
            notes = '\n'.join(notes_parts) if notes_parts else content

            # If no notes, use content as narration
            if not notes and content:
                notes = f"{title}. {content}"
            elif not notes:
                notes = title

            slides.append(SlideContent(
                title=title or f"Slide {slide_num + 1}",
                content=content,
                notes=notes,
                duration_hint=max(len(notes) / 15, 5)  # ~15 chars/second, min 5s
            ))

        return slides

    # ========================================================================
    # AUDIO GENERATION
    # ========================================================================

    def _generate_audio(
        self,
        slides: List[SlideContent],
        output_dir: Path
    ) -> List[str]:
        """Generate audio narration for slides using Azure TTS"""
        audio_files = []

        try:
            import azure.cognitiveservices.speech as speechsdk

            speech_config = speechsdk.SpeechConfig(
                subscription=AZURE_TTS_KEY,
                region=AZURE_TTS_REGION
            )
            speech_config.speech_synthesis_voice_name = AZURE_TTS_VOICE
            speech_config.set_speech_synthesis_output_format(
                speechsdk.SpeechSynthesisOutputFormat.Audio16Khz32KBitRateMonoMp3
            )

            for i, slide in enumerate(slides):
                output_path = output_dir / f"audio_{i}.mp3"

                audio_config = speechsdk.audio.AudioOutputConfig(
                    filename=str(output_path)
                )

                synthesizer = speechsdk.SpeechSynthesizer(
                    speech_config=speech_config,
                    audio_config=audio_config
                )

                # Use SSML for better control
                ssml = f"""
<speak version="1.0" xmlns="http://www.w3.org/2001/10/synthesis" xml:lang="en-US">
    <voice name="{AZURE_TTS_VOICE}">
        <prosody rate="0.9" pitch="0%">
            {slide.notes}
        </prosody>
    </voice>
</speak>
"""

                result = synthesizer.speak_ssml_async(ssml).get()

                if result.reason == speechsdk.ResultReason.SynthesizingAudioCompleted:
                    audio_files.append(str(output_path))
                else:
                    # Fallback: create silent audio
                    audio_files.append(None)

        except ImportError:
            # Azure SDK not available, use gTTS fallback
            try:
                from gtts import gTTS

                for i, slide in enumerate(slides):
                    output_path = output_dir / f"audio_{i}.mp3"
                    tts = gTTS(text=slide.notes, lang='en', slow=False)
                    tts.save(str(output_path))
                    audio_files.append(str(output_path))

            except ImportError:
                # No TTS available, return empty
                audio_files = [None] * len(slides)

        except Exception as e:
            print(f"TTS error: {e}")
            audio_files = [None] * len(slides)

        return audio_files

    # ========================================================================
    # SLIDE RENDERING
    # ========================================================================

    def _render_slides(
        self,
        slides: List[SlideContent],
        output_dir: Path
    ) -> List[str]:
        """Render slides to images"""
        if not PIL_AVAILABLE:
            raise Exception("PIL not available. Install: pip install Pillow")

        image_files = []

        for i, slide in enumerate(slides):
            img = Image.new('RGB', (self.VIDEO_WIDTH, self.VIDEO_HEIGHT), self.BG_COLOR)
            draw = ImageDraw.Draw(img)

            # Load fonts
            try:
                if self.FONT_PATH:
                    title_font = ImageFont.truetype(self.FONT_PATH, self.FONT_TITLE_SIZE)
                    content_font = ImageFont.truetype(self.FONT_PATH, self.FONT_CONTENT_SIZE)
                else:
                    title_font = ImageFont.load_default()
                    content_font = ImageFont.load_default()
            except Exception:
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()

            # Draw accent bar at top
            draw.rectangle(
                [(0, 0), (self.VIDEO_WIDTH, 8)],
                fill=self.ACCENT_COLOR
            )

            # Draw title
            title_y = 100
            self._draw_text_wrapped(
                draw,
                slide.title,
                title_font,
                self.TITLE_COLOR,
                100,
                title_y,
                self.VIDEO_WIDTH - 200
            )

            # Draw content
            content_y = 300
            self._draw_text_wrapped(
                draw,
                slide.content,
                content_font,
                self.CONTENT_COLOR,
                100,
                content_y,
                self.VIDEO_WIDTH - 200
            )

            # Save image
            output_path = output_dir / f"slide_{i}.png"
            img.save(str(output_path))
            image_files.append(str(output_path))

        return image_files

    def _draw_text_wrapped(
        self,
        draw: 'ImageDraw.Draw',
        text: str,
        font: 'ImageFont.FreeTypeFont',
        color: tuple,
        x: int,
        y: int,
        max_width: int
    ):
        """Draw text with word wrapping"""
        lines = text.split('\n')
        current_y = y

        for line in lines:
            words = line.split()
            current_line = ""

            for word in words:
                test_line = f"{current_line} {word}".strip()
                bbox = draw.textbbox((0, 0), test_line, font=font)
                text_width = bbox[2] - bbox[0]

                if text_width <= max_width:
                    current_line = test_line
                else:
                    if current_line:
                        draw.text((x, current_y), current_line, font=font, fill=color)
                        current_y += font.size + 10
                    current_line = word

            if current_line:
                draw.text((x, current_y), current_line, font=font, fill=color)
                current_y += font.size + 20

    # ========================================================================
    # VIDEO CREATION
    # ========================================================================

    def _create_video(
        self,
        slide_images: List[str],
        audio_files: List[str],
        output_path: str
    ) -> float:
        """Combine images and audio into video"""
        if not MOVIEPY_AVAILABLE:
            raise Exception("MoviePy not available. Install: pip install moviepy")

        clips = []
        total_duration = 0

        for i, (img_path, audio_path) in enumerate(zip(slide_images, audio_files)):
            # Determine duration from audio or use default
            if audio_path and Path(audio_path).exists():
                audio_clip = AudioFileClip(audio_path)
                duration = audio_clip.duration + 1  # Add 1 second buffer
            else:
                duration = 5.0  # Default 5 seconds per slide

            # Create image clip
            img_clip = ImageClip(img_path).set_duration(duration)

            # Add audio if available
            if audio_path and Path(audio_path).exists():
                audio_clip = AudioFileClip(audio_path)
                img_clip = img_clip.set_audio(audio_clip)

            clips.append(img_clip)
            total_duration += duration

        # Concatenate all clips
        final_clip = concatenate_videoclips(clips, method="compose")

        # Write output
        final_clip.write_videofile(
            output_path,
            fps=self.FPS,
            codec='libx264',
            audio_codec='aac',
            temp_audiofile='temp-audio.m4a',
            remove_temp=True,
            verbose=False,
            logger=None
        )

        # Cleanup
        final_clip.close()
        for clip in clips:
            clip.close()

        return total_duration

    def _generate_thumbnail(self, first_slide_path: str, output_path: str):
        """Generate video thumbnail"""
        if not PIL_AVAILABLE:
            return

        try:
            img = Image.open(first_slide_path)
            img.thumbnail((480, 270))  # 16:9 thumbnail
            img.save(output_path, "JPEG", quality=85)
        except Exception:
            pass

    # ========================================================================
    # VIDEO MANAGEMENT
    # ========================================================================

    def get_video(self, video_id: str, tenant_id: str) -> Optional[Video]:
        """Get video by ID"""
        return self.db.query(Video).filter(
            Video.id == video_id,
            Video.tenant_id == tenant_id
        ).first()

    def list_videos(
        self,
        tenant_id: str,
        project_id: Optional[str] = None,
        status: Optional[VideoStatus] = None,
        limit: int = 50,
        offset: int = 0
    ) -> Tuple[List[Video], int]:
        """List videos with filtering"""
        query = self.db.query(Video).filter(
            Video.tenant_id == tenant_id
        )

        if project_id:
            query = query.filter(Video.project_id == project_id)
        if status:
            query = query.filter(Video.status == status)

        total = query.count()
        videos = query.order_by(
            Video.created_at.desc()
        ).offset(offset).limit(limit).all()

        return videos, total

    def delete_video(self, video_id: str, tenant_id: str) -> Tuple[bool, Optional[str]]:
        """Delete a video and its files"""
        try:
            video = self.db.query(Video).filter(
                Video.id == video_id,
                Video.tenant_id == tenant_id
            ).first()

            if not video:
                return False, "Video not found"

            # Delete files
            if video.file_path and Path(video.file_path).exists():
                Path(video.file_path).unlink()
            if video.thumbnail_path and Path(video.thumbnail_path).exists():
                Path(video.thumbnail_path).unlink()

            self.db.delete(video)
            self.db.commit()

            return True, None

        except Exception as e:
            self.db.rollback()
            return False, str(e)

    def register_progress_callback(
        self,
        video_id: str,
        callback: Callable[[VideoProgress], None]
    ):
        """Register callback for progress updates"""
        self._progress_callbacks[video_id] = callback

    def unregister_progress_callback(self, video_id: str):
        """Unregister progress callback"""
        if video_id in self._progress_callbacks:
            del self._progress_callbacks[video_id]
