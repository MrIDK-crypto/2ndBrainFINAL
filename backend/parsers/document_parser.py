"""
Document Parser for Office Files
Uses LlamaParse for all document parsing with GPT-4o-mini processing
Falls back to traditional parsers if LlamaParse is unavailable
"""

import os
from pathlib import Path
from typing import Dict, Optional
import warnings
warnings.filterwarnings('ignore')

# Try to import LlamaParse parser
try:
    from parsers.llamaparse_parser import LlamaParseDocumentParser
    HAS_LLAMAPARSE = True
except ImportError:
    HAS_LLAMAPARSE = False

# PDF parsing (fallback)
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

# PowerPoint parsing (fallback)
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Excel parsing (fallback)
try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

# Word parsing (fallback)
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False


class DocumentParser:
    """Parse various document formats to extract text"""

    def __init__(self, config=None, use_llamaparse=True):
        """
        Initialize document parser

        Args:
            config: Configuration object (required if use_llamaparse=True)
            use_llamaparse: Whether to use LlamaParse (default: True)
        """
        self.config = config
        self.use_llamaparse = use_llamaparse and HAS_LLAMAPARSE
        self.llamaparse_parser = None

        # Initialize LlamaParse if available and requested
        if self.use_llamaparse and config:
            try:
                self.llamaparse_parser = LlamaParseDocumentParser(config)
                print("✓ Using LlamaParse for document parsing")
            except Exception as e:
                print(f"⚠ Failed to initialize LlamaParse: {e}")
                print("  Falling back to traditional parsers")
                self.use_llamaparse = False

        # Set up supported formats
        self.supported_formats = []
        if self.use_llamaparse and self.llamaparse_parser:
            self.supported_formats = self.llamaparse_parser.supported_formats
        else:
            if HAS_PDF:
                self.supported_formats.append('.pdf')
            if HAS_PPTX:
                self.supported_formats.append('.pptx')
            if HAS_XLSX:
                self.supported_formats.append('.xlsx')
            if HAS_DOCX:
                self.supported_formats.append('.docx')

    def can_parse(self, file_path: str) -> bool:
        """Check if file format is supported"""
        ext = Path(file_path).suffix.lower()
        return ext in self.supported_formats

    def parse(self, file_path: str) -> Optional[Dict]:
        """
        Parse a document and return extracted text
        Uses LlamaParse if available, falls back to traditional parsers

        Returns:
            Dict with 'content' and 'metadata' or None if parsing failed
        """
        if not os.path.exists(file_path):
            return None

        ext = Path(file_path).suffix.lower()

        # Try LlamaParse first if available
        if self.use_llamaparse and self.llamaparse_parser:
            try:
                return self.llamaparse_parser.parse(file_path)
            except Exception as e:
                print(f"  ⚠ LlamaParse failed: {e}")
                print(f"  Falling back to traditional parser for {Path(file_path).name}")

        # Fall back to traditional parsers
        try:
            if ext == '.pdf' and HAS_PDF:
                return self._parse_pdf(file_path)
            elif ext == '.pptx' and HAS_PPTX:
                return self._parse_pptx(file_path)
            elif ext == '.xlsx' and HAS_XLSX:
                return self._parse_xlsx(file_path)
            elif ext == '.docx' and HAS_DOCX:
                return self._parse_docx(file_path)
        except Exception as e:
            print(f"  ⚠ Error parsing {Path(file_path).name}: {e}")
            return None

        return None

    def _parse_pdf(self, file_path: str) -> Dict:
        """Extract text from PDF"""
        text_parts = []

        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            num_pages = len(pdf_reader.pages)

            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                if text.strip():
                    text_parts.append(text.strip())

        content = '\n\n'.join(text_parts)

        return {
            'content': content,
            'metadata': {
                'pages': num_pages,
                'file_type': 'pdf'
            }
        }

    def _parse_pptx(self, file_path: str) -> Dict:
        """Extract text from PowerPoint"""
        text_parts = []

        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + '\n'.join(slide_text))

        content = '\n\n'.join(text_parts)

        return {
            'content': content,
            'metadata': {
                'slides': len(prs.slides),
                'file_type': 'pptx'
            }
        }

    def _parse_xlsx(self, file_path: str) -> Dict:
        """Extract text from Excel - NO row limit, processes entire spreadsheet"""
        text_parts = []
        total_rows = 0

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"[Sheet: {sheet_name}]"]

            # Read ALL rows - no artificial limit
            # Excel files are typically <50MB so this is fine
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                # Filter out None values and convert to strings
                row_values = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                if row_values:
                    sheet_text.append(' | '.join(row_values))
                    row_count += 1

            if row_count > 0:
                text_parts.append('\n'.join(sheet_text))
                total_rows += row_count

        content = '\n\n'.join(text_parts)

        return {
            'content': content,
            'metadata': {
                'sheets': len(wb.sheetnames),
                'total_rows': total_rows,
                'file_type': 'xlsx'
            }
        }

    def _parse_docx(self, file_path: str) -> Dict:
        """Extract text from Word document"""
        doc = Document(file_path)

        text_parts = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)

        content = '\n\n'.join(text_parts)

        return {
            'content': content,
            'metadata': {
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'file_type': 'docx'
            }
        }

    def parse_pdf_bytes(self, content: bytes) -> str:
        """
        Extract text from PDF bytes (for file uploads)

        Args:
            content: PDF file content as bytes

        Returns:
            Extracted text string
        """
        import io
        text_parts = []

        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            num_pages = len(pdf_reader.pages)

            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                if text.strip():
                    text_parts.append(text.strip())

            return '\n\n'.join(text_parts)
        except Exception as e:
            raise Exception(f"Failed to parse PDF: {str(e)}")

    def parse_word_bytes(self, content: bytes) -> str:
        """
        Extract text from Word document bytes (for file uploads)

        Args:
            content: DOCX file content as bytes

        Returns:
            Extracted text string
        """
        import io
        text_parts = []

        try:
            doc = Document(io.BytesIO(content))

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text.strip())

            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        text_parts.append(row_text)

            return '\n\n'.join(text_parts)
        except Exception as e:
            raise Exception(f"Failed to parse Word document: {str(e)}")


if __name__ == "__main__":
    # Test the parser
    parser = DocumentParser()
    print(f"Supported formats: {parser.supported_formats}")

    # Test with a sample file
    test_file = "/Users/rishitjain/Downloads/Takeout/Google Chat/Groups/Space AAAAn7sv4eE/File-Timeline - BEAT Healthcare Consulting.pptx"
    if os.path.exists(test_file):
        result = parser.parse(test_file)
        if result:
            print(f"\nExtracted {len(result['content'])} characters")
            print(f"Preview: {result['content'][:200]}...")
