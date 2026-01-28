"""
OneDrive Connector
Connects to Microsoft OneDrive/SharePoint to extract PowerPoint, Excel, and Word files.
"""

import os
import io
import mimetypes
from datetime import datetime
from typing import List, Dict, Optional, Any

from .base_connector import BaseConnector, ConnectorConfig, ConnectorStatus, Document

# Note: Requires msal and requests
# pip install msal requests

try:
    import msal
    import requests
    ONEDRIVE_AVAILABLE = True
except ImportError:
    ONEDRIVE_AVAILABLE = False


class OneDriveConnector(BaseConnector):
    """
    OneDrive connector for extracting Microsoft Office files.

    Extracts:
    - PowerPoint presentations (.pptx, .ppt)
    - Excel spreadsheets (.xlsx, .xls)
    - Word documents (.docx, .doc)
    - PDFs
    """

    CONNECTOR_TYPE = "onedrive"
    REQUIRED_CREDENTIALS = ["access_token", "refresh_token"]
    OPTIONAL_SETTINGS = {
        "folder_ids": [],  # Specific folders to sync (empty = root)
        "file_types": [".pptx", ".ppt", ".xlsx", ".xls", ".docx", ".doc", ".pdf"],
        "max_file_size_mb": 50,
        "include_shared": True
    }

    # Microsoft Graph API endpoint
    GRAPH_ENDPOINT = "https://graph.microsoft.com/v1.0"

    def __init__(self, config: ConnectorConfig):
        super().__init__(config)
        self.access_token = None

    @classmethod
    def _get_oauth_config(cls) -> Dict:
        return {
            "client_id": os.getenv("MICROSOFT_CLIENT_ID", ""),
            "client_secret": os.getenv("MICROSOFT_CLIENT_SECRET", ""),
            "tenant": os.getenv("MICROSOFT_TENANT_ID", "common"),
            "redirect_uri": os.getenv("MICROSOFT_REDIRECT_URI", "http://localhost:5003/api/integrations/onedrive/callback")
        }

    @classmethod
    def get_auth_url(cls, redirect_uri: str, state: str) -> str:
        """Generate Microsoft OAuth authorization URL"""
        if not ONEDRIVE_AVAILABLE:
            raise ImportError("MSAL not installed")

        config = cls._get_oauth_config()

        # Create MSAL app (ConfidentialClient for server-side apps with client_secret)
        app = msal.ConfidentialClientApplication(
            config["client_id"],
            authority=f"https://login.microsoftonline.com/{config['tenant']}",
            client_credential=config["client_secret"]
        )

        # Generate auth URL
        scopes = ["Files.Read.All", "User.Read"]

        auth_url = app.get_authorization_request_url(
            scopes,
            state=state,
            redirect_uri=redirect_uri
        )

        return auth_url

    @classmethod
    def exchange_code_for_tokens(cls, code: str, redirect_uri: str):
        """Exchange authorization code for access token"""
        try:
            if not ONEDRIVE_AVAILABLE:
                return None, "MSAL not installed. Run: pip install msal"

            config = cls._get_oauth_config()

            # Create MSAL app (ConfidentialClient for server-side apps with client_secret)
            app = msal.ConfidentialClientApplication(
                config["client_id"],
                authority=f"https://login.microsoftonline.com/{config['tenant']}",
                client_credential=config["client_secret"]
            )

            # Exchange code for token
            result = app.acquire_token_by_authorization_code(
                code,
                scopes=["Files.Read.All", "User.Read"],
                redirect_uri=redirect_uri
            )

            if "error" in result:
                return None, result.get("error_description", result["error"])

            tokens = {
                "access_token": result["access_token"],
                "refresh_token": result.get("refresh_token"),
                "expires_in": result.get("expires_in")
            }

            return tokens, None

        except Exception as e:
            return None, str(e)

    async def connect(self) -> bool:
        """Connect to OneDrive"""
        if not ONEDRIVE_AVAILABLE:
            self._set_error("MSAL not installed. Run: pip install msal")
            return False

        try:
            self.status = ConnectorStatus.CONNECTING
            self.access_token = self.config.credentials.get("access_token")

            # Test connection by getting user profile
            response = requests.get(
                f"{self.GRAPH_ENDPOINT}/me",
                headers={"Authorization": f"Bearer {self.access_token}"}
            )

            if response.status_code != 200:
                self._set_error(f"Connection failed: {response.text}")
                return False

            user_data = response.json()
            self.sync_stats["user"] = user_data.get("displayName", "Unknown")

            self.status = ConnectorStatus.CONNECTED
            self._clear_error()
            return True

        except Exception as e:
            self._set_error(f"Failed to connect: {str(e)}")
            return False

    async def disconnect(self) -> bool:
        """Disconnect from OneDrive"""
        self.access_token = None
        self.status = ConnectorStatus.DISCONNECTED
        return True

    async def test_connection(self) -> bool:
        """Test OneDrive connection"""
        if not self.access_token:
            return False

        try:
            response = requests.get(
                f"{self.GRAPH_ENDPOINT}/me/drive",
                headers={"Authorization": f"Bearer {self.access_token}"}
            )
            return response.status_code == 200
        except Exception:
            return False

    async def sync(self, since: Optional[datetime] = None) -> List[Document]:
        """Sync files from OneDrive"""
        if not self.access_token:
            await self.connect()

        if self.status != ConnectorStatus.CONNECTED:
            return []

        self.status = ConnectorStatus.SYNCING
        documents = []

        try:
            # Get folders to sync
            folder_ids = self.config.settings.get("folder_ids", [])

            if folder_ids:
                # Sync specific folders
                for folder_id in folder_ids:
                    folder_docs = await self._sync_folder(folder_id, since)
                    documents.extend(folder_docs)
            else:
                # Sync root folder
                folder_docs = await self._sync_folder("root", since)
                documents.extend(folder_docs)

            # Update stats
            self.sync_stats = {
                "documents_synced": len(documents),
                "sync_time": datetime.now().isoformat()
            }

            self.config.last_sync = datetime.now()
            self.status = ConnectorStatus.CONNECTED

        except Exception as e:
            self._set_error(f"Sync failed: {str(e)}")

        return documents

    async def _sync_folder(self, folder_id: str, since: Optional[datetime]) -> List[Document]:
        """Sync files from a folder recursively"""
        documents = []

        try:
            # Get files in folder
            url = f"{self.GRAPH_ENDPOINT}/me/drive/items/{folder_id}/children"
            response = requests.get(
                url,
                headers={"Authorization": f"Bearer {self.access_token}"}
            )

            if response.status_code != 200:
                print(f"[OneDrive] Failed to list folder {folder_id}: {response.text}")
                return documents

            items = response.json().get("value", [])

            for item in items:
                # Check if it's a folder
                if "folder" in item:
                    # Recursively sync subfolder
                    subfolder_docs = await self._sync_folder(item["id"], since)
                    documents.extend(subfolder_docs)

                # Check if it's a file
                elif "file" in item:
                    # Check file type
                    name = item.get("name", "")
                    file_types = self.config.settings.get("file_types", [])

                    if any(name.lower().endswith(ext) for ext in file_types):
                        # Check file size
                        size_mb = item.get("size", 0) / (1024 * 1024)
                        max_size = self.config.settings.get("max_file_size_mb", 50)

                        if size_mb > max_size:
                            print(f"[OneDrive] Skipping {name} - too large ({size_mb:.1f}MB)")
                            continue

                        # Check modification time
                        modified = datetime.fromisoformat(item["lastModifiedDateTime"].replace("Z", "+00:00"))

                        if since and modified < since:
                            continue

                        # Download and parse file
                        doc = await self._download_and_parse(item)
                        if doc:
                            documents.append(doc)

        except Exception as e:
            print(f"[OneDrive] Error syncing folder {folder_id}: {e}")

        return documents

    async def _download_and_parse(self, item: Dict) -> Optional[Document]:
        """Download and parse a file"""
        try:
            file_id = item["id"]
            name = item.get("name", "unknown")
            download_url = item.get("@microsoft.graph.downloadUrl")

            if not download_url:
                print(f"[OneDrive] No download URL for {name}")
                return None

            print(f"[OneDrive] Downloading {name}...")

            # Download file
            response = requests.get(download_url)

            if response.status_code != 200:
                print(f"[OneDrive] Failed to download {name}")
                return None

            file_content = response.content

            # Parse based on file type
            content_text = await self._parse_file(name, file_content)

            if not content_text:
                return None

            # Create document
            modified = datetime.fromisoformat(item["lastModifiedDateTime"].replace("Z", "+00:00"))

            return Document(
                doc_id=f"onedrive_{file_id}",
                source="onedrive",
                content=content_text,
                title=name,
                metadata={
                    "file_id": file_id,
                    "size": item.get("size", 0),
                    "path": item.get("parentReference", {}).get("path", ""),
                    "web_url": item.get("webUrl")
                },
                timestamp=modified,
                author=item.get("createdBy", {}).get("user", {}).get("displayName"),
                url=item.get("webUrl"),
                doc_type=self._get_doc_type(name)
            )

        except Exception as e:
            print(f"[OneDrive] Error parsing {item.get('name')}: {e}")
            return None

    async def _parse_file(self, filename: str, content: bytes) -> Optional[str]:
        """Parse file content based on type"""
        try:
            lower_name = filename.lower()

            # PowerPoint
            if lower_name.endswith((".pptx", ".ppt")):
                return self._parse_powerpoint(content)

            # Excel
            elif lower_name.endswith((".xlsx", ".xls")):
                return self._parse_excel(content)

            # Word
            elif lower_name.endswith((".docx", ".doc")):
                return self._parse_word(content)

            # PDF
            elif lower_name.endswith(".pdf"):
                return self._parse_pdf(content)

            return None

        except Exception as e:
            print(f"[OneDrive] Parse error for {filename}: {e}")
            return None

    def _parse_powerpoint(self, content: bytes) -> Optional[str]:
        """Parse PowerPoint file"""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            text_parts = []

            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- Slide {i} ---\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"

                text_parts.append(slide_text)

            return "\n".join(text_parts)

        except Exception as e:
            print(f"[OneDrive] PowerPoint parse error: {e}")
            return None

    def _parse_excel(self, content: bytes) -> Optional[str]:
        """Parse Excel file"""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n--- Sheet: {sheet_name} ---\n")

                # Get used range
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)

            return "\n".join(text_parts)

        except Exception as e:
            print(f"[OneDrive] Excel parse error: {e}")
            return None

    def _parse_word(self, content: bytes) -> Optional[str]:
        """Parse Word document"""
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

            return "\n\n".join(paragraphs)

        except Exception as e:
            print(f"[OneDrive] Word parse error: {e}")
            return None

    def _parse_pdf(self, content: bytes) -> Optional[str]:
        """Parse PDF file"""
        try:
            import PyPDF2

            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            text_parts = []

            for page in pdf_reader.pages:
                text_parts.append(page.extract_text())

            return "\n\n".join(text_parts)

        except Exception as e:
            print(f"[OneDrive] PDF parse error: {e}")
            return None

    def _get_doc_type(self, filename: str) -> str:
        """Get document type based on filename"""
        lower_name = filename.lower()

        if lower_name.endswith((".pptx", ".ppt")):
            return "presentation"
        elif lower_name.endswith((".xlsx", ".xls")):
            return "spreadsheet"
        elif lower_name.endswith((".docx", ".doc")):
            return "document"
        elif lower_name.endswith(".pdf"):
            return "pdf"

        return "file"

    async def get_document(self, doc_id: str) -> Optional[Document]:
        """Get a specific document by ID"""
        # Implementation would fetch single file by ID
        return None
